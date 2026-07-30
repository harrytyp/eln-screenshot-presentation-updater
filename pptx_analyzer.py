"""
pptx_analyzer.py — Extract images from a PPTX and classify them.

For each image in the presentation, this module determines:
- Whether it is a "real screenshot" (should be replaced) or an icon/logo/diagram
- Which ELN page it likely comes from (based on slide context + mapping)
- The exact position/size it occupies in the slide (for replacement)
"""

from __future__ import annotations

import io
import os
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional

import yaml
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu


@dataclass
class ImageInfo:
    """Describes an image embedded in a PPTX slide."""
    slide_index: int          # 0-based
    slide_number: int         # 1-based
    shape_name: str           # e.g. "Grafik 19"
    blob: bytes
    content_type: str         # e.g. "image/png"
    width_emu: int            # displayed width in EMU
    height_emu: int           # displayed height in EMU
    left_emu: int             # position left in EMU
    top_emu: int              # position top in EMU
    width_px: int             # actual pixel width
    height_px: int            # actual pixel height
    is_likely_screenshot: bool = False
    page_url: str = ""        # Identified/mapped ELN URL
    page_confidence: str = "" # "high", "medium", "low", or ""
    slide_context: str = ""   # Text from the same slide for context


# Thresholds for classifying images
SCREENSHOT_MIN_SIZE = 30_000       # 30KB minimum for a real screenshot
SCREENSHOT_MIN_DIMENSION = 200     # 200px minimum width or height
ICON_MAX_SIZE = 20_000             # Under 20KB is likely an icon
NON_SCREENSHOT_PATTERNS = [
    r"(?i)logo",
    r"(?i)qrcode",
    r"(?i)qr.code",
]


def extract_images(pptx_path: str, output_dir: Optional[str] = None) -> list[ImageInfo]:
    """
    Extract all images from a PPTX file.
    Optionally save them to output_dir.
    Returns a list of ImageInfo records.
    """
    prs = Presentation(pptx_path)
    images: list[ImageInfo] = []

    if output_dir:
        os.makedirs(output_dir, exist_ok=True)

    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1

        # Collect slide text for context
        slide_text_parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    slide_text_parts.append(t)
        slide_context = " | ".join(slide_text_parts[:5])

        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue

            img = shape.image
            blob = img.blob
            content_type = img.content_type

            # Get displayed dimensions
            w_emu = shape.width
            h_emu = shape.height
            l_emu = shape.left
            t_emu = shape.top

            # Get actual pixel dimensions from the image
            pil_img = Image.open(io.BytesIO(blob))
            px_w, px_h = pil_img.size

            info = ImageInfo(
                slide_index=slide_idx,
                slide_number=slide_num,
                shape_name=shape.name,
                blob=blob,
                content_type=content_type,
                width_emu=w_emu,
                height_emu=h_emu,
                left_emu=l_emu,
                top_emu=t_emu,
                width_px=px_w,
                height_px=px_h,
            )

            # Classify
            _classify_image(info, slide_context)

            images.append(info)

            # Save if requested
            if output_dir:
                ext = content_type.split("/")[-1]
                if ext == "jpeg":
                    ext = "jpg"
                safe_name = shape.name.replace(" ", "_")
                fname = f"s{slide_num:02d}_{safe_name}.{ext}"
                fpath = os.path.join(output_dir, fname)
                with open(fpath, "wb") as f:
                    f.write(blob)

    return images


def _classify_image(info: ImageInfo, slide_context: str):
    """Determine if this image is a screenshot, and which page it shows."""
    info.slide_context = slide_context

    # Rule 1: Too small -> not a screenshot
    if len(info.blob) < ICON_MAX_SIZE:
        info.is_likely_screenshot = False
        info.page_confidence = ""
        return

    # Rule 2: Dimensions too small -> not a screenshot
    if info.width_px < SCREENSHOT_MIN_DIMENSION or info.height_px < SCREENSHOT_MIN_DIMENSION:
        info.is_likely_screenshot = False
        info.page_confidence = ""
        return

    # Rule 3: Check shape name patterns for exclusions
    for pat in NON_SCREENSHOT_PATTERNS:
        if re.search(pat, info.shape_name):
            info.is_likely_screenshot = False
            info.page_confidence = ""
            return

    # Rule 4: Very large images (>3000px) that aren't screenshots of the web UI
    # (likely full-page diagrams or offboarding charts)
    if info.width_px > 3000 or info.height_px > 2000:
        # Check if the context mentions ELN, eLabFTW etc.
        eln_keywords = ["elabftw", "screenshot", "dashboard", "experiment", "resource",
                        "scheduler", "template", "permission", "data"]
        if not any(kw in slide_context.lower() for kw in eln_keywords):
            info.is_likely_screenshot = False
            info.page_confidence = "low"
            return

    # If it passed the size checks, it's likely a screenshot worth updating
    info.is_likely_screenshot = True


def map_screenshots(images: list[ImageInfo], mapping_path: str) -> list[ImageInfo]:
    """
    Load the screenshot mapping YAML and match each image to a known ELN page.
    Images that don't have a mapping entry are flagged as unmapped.
    """
    with open(mapping_path, "r", encoding="utf-8") as f:
        mapping_data = yaml.safe_load(f)

    entries = mapping_data.get("screenshots", [])
    unmapped: list[ImageInfo] = []

    for img_info in images:
        if not img_info.is_likely_screenshot:
            continue

        # Find matching entry in mapping
        matched = False
        for entry in entries:
            if (entry["slide"] == img_info.slide_number and
                entry["shape"] == img_info.shape_name):
                img_info.page_url = entry["url"]
                img_info.page_confidence = "high"
                # Store additional mapping metadata on the object
                img_info._mapping = entry  # type: ignore[attr-defined]
                matched = True
                break

        if not matched:
            img_info.page_confidence = "unmapped"
            unmapped.append(img_info)

    return unmapped


def generate_report(images: list[ImageInfo]) -> str:
    """Generate a human-readable report of all images and their mapping."""
    lines = []
    lines.append("=" * 70)
    lines.append("PPTX SCREENSHOT ANALYSIS REPORT")
    lines.append("=" * 70)
    lines.append("")

    screenshots = [i for i in images if i.is_likely_screenshot]
    icons = [i for i in images if not i.is_likely_screenshot]

    lines.append(f"Total images: {len(images)}")
    lines.append(f"Likely screenshots (will be updated): {len(screenshots)}")
    lines.append(f"Icons/logos/diagrams (left unchanged):  {len(icons)}")
    lines.append("")

    if screenshots:
        lines.append("-- SCREENSHOTS TO UPDATE --")
        lines.append(f"{'Slide':>5} {'Shape':<20} {'Dimensions':<14} {'Size':<10} {'Page URL':<30}")
        lines.append("-" * 85)
        for img in sorted(screenshots, key=lambda x: x.slide_number):
            size_kb = len(img.blob) / 1024
            url = getattr(img, 'page_url', '')
            conf = getattr(img, 'page_confidence', '')
            page = f"{url} [{conf}]" if conf else "UNMAPPED"
            dims = f"{img.width_px}x{img.height_px}"
            lines.append(f"{img.slide_number:>5} {img.shape_name:<20} {dims:<14} {size_kb:>6.1f}KB  {page:<30}")

    if icons:
        lines.append("")
        lines.append("-- ICONS / DIAGRAMS (NOT updated) --")
        for img in sorted(icons, key=lambda x: x.slide_number):
            size_kb = len(img.blob) / 1024
            lines.append(f"  Slide {img.slide_number:>2} {img.shape_name:<25} {img.width_px}x{img.height_px:<6} {size_kb:>5.1f}KB")

    return "\n".join(lines)


def load_env(env_path: str = ".env") -> dict:
    """Load environment variables from .env file."""
    result = {}
    if not os.path.exists(env_path):
        return result
    with open(env_path, "r") as f:
        for line in f:
            line = line.strip()
            if not line or line.startswith("#"):
                continue
            if "=" in line:
                key, _, val = line.partition("=")
                result[key.strip()] = val.strip()
    return result


# Re-export ImageInfo for use by other modules (already top-level)
