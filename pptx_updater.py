"""
pptx_updater.py — Replace images in a PPTX with new screenshots.

Works by parsing the PPTX XML to find image relationships and swapping
the image blobs while preserving position, size, and other shape properties.
"""

from __future__ import annotations

import io
import os
import shutil
from pathlib import Path
from typing import Optional

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from pptx_analyzer import ImageInfo, extract_images


def replace_screenshot_images(
    pptx_path: str,
    output_path: str,
    captured: dict[str, str],
    mapping_path: str,
) -> tuple[int, list[str]]:
    """
    Replace old screenshots in a PPTX with newly captured ones.

    Args:
        pptx_path: Path to the original PPTX
        output_path: Path to save the updated PPTX
        captured: Dict mapping "s{slide}_{shape_name}" -> path to new screenshot
        mapping_path: Path to screenshot mapping YAML (for verifying which to update)

    Returns:
        (number_of_replacements, list_of_warnings)
    """
    from pptx_analyzer import map_screenshots

    warnings: list[str] = []

    # Extract images from the original PPTX
    images = extract_images(pptx_path)
    unmapped = map_screenshots(images, mapping_path)

    if unmapped:
        for img in unmapped:
            warnings.append(
                f"Slide {img.slide_number} '{img.shape_name}' — "
                f"no mapping entry; screenshot NOT replaced."
            )

    # Open the PPTX for modification
    prs = Presentation(pptx_path)

    replacements = 0
    for image_info in images:
        if not image_info.is_likely_screenshot:
            continue

        # Build key in the same format as captured dict
        key = f"s{image_info.slide_number:02d}_{image_info.shape_name}"

        if key not in captured:
            warnings.append(f"Slide {image_info.slide_number} '{image_info.shape_name}' — no captured screenshot available")
            continue

        new_image_path = captured[key]
        if not os.path.exists(new_image_path):
            warnings.append(f"Captured file not found: {new_image_path}")
            continue

        # Find the matching shape in the slide and replace image data
        slide = prs.slides[image_info.slide_index]
        replaced = _replace_shape_image(slide, image_info.shape_name, new_image_path)
        if replaced:
            replacements += 1
        else:
            warnings.append(f"Slide {image_info.slide_number} '{image_info.shape_name}' — shape not found for replacement")

    # Save the updated PPTX
    prs.save(output_path)
    return replacements, warnings


def _replace_shape_image(slide, shape_name: str, new_image_path: str) -> bool:
    """
    Find a picture shape by name in the slide and replace its image blob.
    Uses rel.target_part.blob to persist the replacement through save.
    Returns True on success.
    """
    NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        if shape.name != shape_name:
            continue

        # Read new image
        with open(new_image_path, "rb") as f:
            new_blob = f.read()

        if not new_blob:
            print(f"  WARNING: Empty blob for {new_image_path}")
            return False

        # Find the blip element and its rId
        blip = shape._element.find(f".//{{{NS_A}}}blip")
        if blip is None:
            return False

        rId = blip.get(f"{{{NS_R}}}embed")
        if not rId:
            return False

        # Replace the blob on the image part via the relationship
        rel = slide.part.rels[rId]
        rel.target_part.blob = new_blob

        return True

    return False


def resize_screenshots(
    captured: dict[str, str],
    reference_dir: str,
    output_dir: str = "screenshots_resized",
) -> dict[str, str]:
    """
    Resize new screenshots to match the dimensions of the originals.
    This ensures they fit perfectly in their PPTX placeholders.
    """
    os.makedirs(output_dir, exist_ok=True)
    resized: dict[str, str] = {}

    for key, new_path in captured.items():
        # Find the matching original
        # The key format is s{slide}_{shape}
        # We need to find the original in reference_dir
        ref_pattern = key.replace(" ", "_")
        ref_file = None
        for f in os.listdir(reference_dir):
            if f.startswith(ref_pattern) or f.startswith(key.replace(" ", "_")):
                ref_file = os.path.join(reference_dir, f)
                break

        if not ref_file:
            # Try to find by slide number + shape name in any format
            slide_part = key.split("_")[0]  # e.g. "s10"
            for f in os.listdir(reference_dir):
                if f.startswith(slide_part) and key.split("_", 1)[1].replace(" ", "_") in f.replace(" ", "_"):
                    ref_file = os.path.join(reference_dir, f)
                    break

        if not ref_file:
            print(f"  WARNING: No reference image found for {key}, keeping original size")
            resized[key] = new_path
            continue

        # Get reference dimensions
        ref_img = Image.open(ref_file)
        ref_w, ref_h = ref_img.size

        # Open new screenshot and resize
        new_img = Image.open(new_path)
        new_w, new_h = new_img.size

        # Only resize if dimensions differ
        if (new_w, new_h) != (ref_w, ref_h):
            resized_img = new_img.resize((ref_w, ref_h), Image.Resampling.LANCZOS)
            # Convert RGBA to RGB if saving as JPEG
            out_path = os.path.join(output_dir, os.path.basename(new_path))
            if resized_img.mode == "RGBA" and out_path.lower().endswith(".jpg"):
                resized_img = resized_img.convert("RGB")
            resized_img.save(out_path)
            resized[key] = out_path
            print(f"  Resized {os.path.basename(new_path)}: {new_w}x{new_h} → {ref_w}x{ref_h}")
        else:
            # Same dimensions, keep as-is
            resized[key] = new_path

    return resized
